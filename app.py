from flask import Flask, render_template, request, jsonify, session, Response, stream_with_context, redirect, url_for, flash
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
import time
from crewai import Agent, Task, Crew
from crewai.tools import BaseTool
from pydantic import Field
from typing import Type, Any
from pydantic import BaseModel
from dotenv import load_dotenv
import os
import re
import json
import hashlib
from datetime import datetime
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
import base64
from openai import OpenAI
import requests

# Load environment variables
load_dotenv()

app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', os.urandom(24))
app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL', 'sqlite:///ai_assistant.db')
# Fix for PostgreSQL URL from some providers
if app.config['SQLALCHEMY_DATABASE_URI'].startswith('postgres://'):
    app.config['SQLALCHEMY_DATABASE_URI'] = app.config['SQLALCHEMY_DATABASE_URI'].replace('postgres://', 'postgresql://', 1)
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['ALLOWED_EXTENSIONS'] = {
    'txt', 'md', 'pdf',  # Text and PDF
    'doc', 'docx',  # Word
    'ppt', 'pptx',  # PowerPoint
    'xls', 'xlsx',  # Excel
    'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'  # Images (with Vision AI support)
}

db = SQLAlchemy(app)
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

# Create uploads directory if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Database Models
class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(255), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    chat_sessions = db.relationship('ChatSession', backref='user', lazy=True, cascade='all, delete-orphan')
    documents = db.relationship('UploadedDocument', backref='user', lazy=True, cascade='all, delete-orphan')

    def set_password(self, password):
        self.password_hash = generate_password_hash(password)

    def check_password(self, password):
        return check_password_hash(self.password_hash, password)

class ChatSession(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    messages = db.relationship('Message', backref='session', lazy=True, cascade='all, delete-orphan')

class Message(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    session_id = db.Column(db.Integer, db.ForeignKey('chat_session.id'), nullable=False)
    role = db.Column(db.String(20), nullable=False)  # 'user', 'assistant', or 'system'
    content = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

class UploadedDocument(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    filename = db.Column(db.String(255), nullable=False)
    file_hash = db.Column(db.String(64), unique=True, nullable=False)  # MD5 hash
    content = db.Column(db.Text, nullable=False)
    file_type = db.Column(db.String(20), nullable=False)
    uploaded_at = db.Column(db.DateTime, default=datetime.utcnow)

class ExecutionLog(db.Model):
    """Track tool executions for debugging and monitoring"""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)
    tool_name = db.Column(db.String(100), nullable=False)
    tool_type = db.Column(db.String(50), nullable=False)  # 'search', 'vision', 'document', 'agent'
    input_data = db.Column(db.Text, nullable=True)
    output_data = db.Column(db.Text, nullable=True)
    execution_time_ms = db.Column(db.Integer, nullable=True)
    status = db.Column(db.String(20), nullable=False)  # 'success', 'error', 'timeout'
    error_message = db.Column(db.Text, nullable=True)
    extra_info = db.Column(db.Text, nullable=True)  # JSON string for additional info

# Knowledge Graph Models
class KnowledgeEntity(db.Model):
    """Store entities in the knowledge graph (people, topics, preferences, facts)"""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    entity_type = db.Column(db.String(50), nullable=False)  # 'person', 'topic', 'preference', 'fact', 'skill', 'goal', 'location'
    name = db.Column(db.String(255), nullable=False)
    description = db.Column(db.Text, nullable=True)
    attributes = db.Column(db.Text, nullable=True)  # JSON string for additional attributes
    confidence = db.Column(db.Float, default=1.0)  # Confidence score 0-1
    source = db.Column(db.String(100), nullable=True)  # Where this was learned from
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    last_accessed = db.Column(db.DateTime, default=datetime.utcnow)
    access_count = db.Column(db.Integer, default=0)

    # Relationships
    outgoing_relations = db.relationship('KnowledgeRelation', foreign_keys='KnowledgeRelation.source_id', backref='source_entity', lazy=True, cascade='all, delete-orphan')
    incoming_relations = db.relationship('KnowledgeRelation', foreign_keys='KnowledgeRelation.target_id', backref='target_entity', lazy=True, cascade='all, delete-orphan')

class KnowledgeRelation(db.Model):
    """Store relationships between entities in the knowledge graph"""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    source_id = db.Column(db.Integer, db.ForeignKey('knowledge_entity.id'), nullable=False)
    target_id = db.Column(db.Integer, db.ForeignKey('knowledge_entity.id'), nullable=False)
    relation_type = db.Column(db.String(100), nullable=False)  # 'works_at', 'interested_in', 'knows', 'lives_in', 'prefers', etc.
    description = db.Column(db.Text, nullable=True)
    strength = db.Column(db.Float, default=1.0)  # Relationship strength 0-1
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

# Global list for in-memory logs (for tools that run before request context)
execution_logs_memory = []

def log_execution(tool_name, tool_type, input_data, output_data, execution_time_ms, status, error_message=None, extra_info=None, user_id=None):
    """Log a tool execution to the database"""
    try:
        # Truncate long outputs for storage
        input_str = str(input_data)[:5000] if input_data else None
        output_str = str(output_data)[:10000] if output_data else None
        extra_info_str = json.dumps(extra_info) if extra_info else None

        log_entry = ExecutionLog(
            user_id=user_id,
            tool_name=tool_name,
            tool_type=tool_type,
            input_data=input_str,
            output_data=output_str,
            execution_time_ms=execution_time_ms,
            status=status,
            error_message=error_message,
            extra_info=extra_info_str
        )
        db.session.add(log_entry)
        db.session.commit()
    except Exception as e:
        # If DB logging fails, store in memory
        execution_logs_memory.append({
            'timestamp': datetime.utcnow().isoformat(),
            'tool_name': tool_name,
            'tool_type': tool_type,
            'input_data': str(input_data)[:500] if input_data else None,
            'output_data': str(output_data)[:500] if output_data else None,
            'execution_time_ms': execution_time_ms,
            'status': status,
            'error_message': error_message
        })
        print(f"[LOG] Failed to save to DB, stored in memory: {str(e)}")

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# Initialize database
with app.app_context():
    db.create_all()

# Knowledge Graph Functions
def extract_entities_from_conversation(user_message: str, assistant_response: str, user_id: int) -> dict:
    """Use OpenAI to extract entities and relationships from a conversation."""
    try:
        openai_api_key = os.environ.get('OPENAI_API_KEY', '')
        if not openai_api_key:
            return {"entities": [], "relations": []}

        client = OpenAI(api_key=openai_api_key)

        extraction_prompt = f"""Analyze this conversation and extract knowledge about the USER (not general facts).
Extract ONLY information that the user reveals about themselves, their preferences, interests, work, relationships, etc.

User message: {user_message}
Assistant response: {assistant_response}

Return a JSON object with:
{{
    "entities": [
        {{
            "type": "person|topic|preference|fact|skill|goal|location|organization",
            "name": "entity name",
            "description": "brief description",
            "attributes": {{"key": "value"}}
        }}
    ],
    "relations": [
        {{
            "source": "User",
            "relation": "works_at|interested_in|knows|lives_in|prefers|has_skill|wants|dislikes|uses|studies",
            "target": "entity name",
            "description": "optional context"
        }}
    ]
}}

Rules:
- Only extract information the USER explicitly shares about themselves
- "User" should always be the source for relations about the user
- Skip generic conversation that doesn't reveal user information
- Be conservative - only extract clear, explicit information
- Return empty arrays if no personal information is found

Return ONLY valid JSON, no other text."""

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": extraction_prompt}],
            temperature=0.1,
            max_tokens=1000
        )

        result_text = response.choices[0].message.content.strip()
        # Clean up the response (remove markdown code blocks if present)
        if result_text.startswith("```"):
            result_text = result_text.split("```")[1]
            if result_text.startswith("json"):
                result_text = result_text[4:]
        result_text = result_text.strip()

        return json.loads(result_text)
    except Exception as e:
        print(f"[KNOWLEDGE] Entity extraction failed: {str(e)}")
        return {"entities": [], "relations": []}

def store_entity(user_id: int, entity_type: str, name: str, description: str = None, attributes: dict = None, source: str = "conversation") -> KnowledgeEntity:
    """Store or update an entity in the knowledge graph."""
    try:
        # Check if entity already exists for this user
        existing = KnowledgeEntity.query.filter_by(
            user_id=user_id,
            entity_type=entity_type,
            name=name
        ).first()

        if existing:
            # Update existing entity
            if description:
                existing.description = description
            if attributes:
                # Merge attributes
                existing_attrs = json.loads(existing.attributes) if existing.attributes else {}
                existing_attrs.update(attributes)
                existing.attributes = json.dumps(existing_attrs)
            existing.access_count += 1
            existing.last_accessed = datetime.utcnow()
            db.session.commit()
            return existing
        else:
            # Create new entity
            entity = KnowledgeEntity(
                user_id=user_id,
                entity_type=entity_type,
                name=name,
                description=description,
                attributes=json.dumps(attributes) if attributes else None,
                source=source
            )
            db.session.add(entity)
            db.session.commit()
            return entity
    except Exception as e:
        print(f"[KNOWLEDGE] Failed to store entity: {str(e)}")
        db.session.rollback()
        return None

def store_relation(user_id: int, source_entity: KnowledgeEntity, target_entity: KnowledgeEntity, relation_type: str, description: str = None) -> KnowledgeRelation:
    """Store a relationship between two entities."""
    try:
        # Check if relation already exists
        existing = KnowledgeRelation.query.filter_by(
            user_id=user_id,
            source_id=source_entity.id,
            target_id=target_entity.id,
            relation_type=relation_type
        ).first()

        if existing:
            if description:
                existing.description = description
            existing.strength = min(existing.strength + 0.1, 1.0)  # Increase strength
            db.session.commit()
            return existing
        else:
            relation = KnowledgeRelation(
                user_id=user_id,
                source_id=source_entity.id,
                target_id=target_entity.id,
                relation_type=relation_type,
                description=description
            )
            db.session.add(relation)
            db.session.commit()
            return relation
    except Exception as e:
        print(f"[KNOWLEDGE] Failed to store relation: {str(e)}")
        db.session.rollback()
        return None

def update_knowledge_graph(user_id: int, user_message: str, assistant_response: str):
    """Update the knowledge graph based on a conversation exchange."""
    try:
        # Extract entities and relations
        extracted = extract_entities_from_conversation(user_message, assistant_response, user_id)

        if not extracted.get("entities") and not extracted.get("relations"):
            return

        # Ensure "User" entity exists
        user_entity = store_entity(user_id, "person", "User", "The user of this assistant", source="system")

        # Store extracted entities
        entity_map = {"User": user_entity}
        for ent in extracted.get("entities", []):
            stored = store_entity(
                user_id=user_id,
                entity_type=ent.get("type", "topic"),
                name=ent.get("name"),
                description=ent.get("description"),
                attributes=ent.get("attributes"),
                source="conversation"
            )
            if stored:
                entity_map[ent.get("name")] = stored

        # Store relations
        for rel in extracted.get("relations", []):
            source_name = rel.get("source", "User")
            target_name = rel.get("target")

            # Get or create source entity
            if source_name not in entity_map:
                source_ent = store_entity(user_id, "person", source_name, source="conversation")
                entity_map[source_name] = source_ent
            else:
                source_ent = entity_map[source_name]

            # Get or create target entity
            if target_name not in entity_map:
                target_ent = store_entity(user_id, "topic", target_name, source="conversation")
                entity_map[target_name] = target_ent
            else:
                target_ent = entity_map[target_name]

            if source_ent and target_ent:
                store_relation(
                    user_id=user_id,
                    source_entity=source_ent,
                    target_entity=target_ent,
                    relation_type=rel.get("relation", "related_to"),
                    description=rel.get("description")
                )

        print(f"[KNOWLEDGE] Updated graph: {len(extracted.get('entities', []))} entities, {len(extracted.get('relations', []))} relations")
    except Exception as e:
        print(f"[KNOWLEDGE] Failed to update knowledge graph: {str(e)}")

def get_user_knowledge_context(user_id: int, query: str = None, limit: int = 20) -> str:
    """Get relevant knowledge about the user to provide context for responses."""
    try:
        # Get user's entities and relations
        entities = KnowledgeEntity.query.filter_by(user_id=user_id).order_by(
            KnowledgeEntity.access_count.desc(),
            KnowledgeEntity.updated_at.desc()
        ).limit(limit).all()

        if not entities:
            return ""

        context_parts = ["## What I know about you:"]

        # Group entities by type
        by_type = {}
        for ent in entities:
            if ent.name == "User":
                continue
            etype = ent.entity_type
            if etype not in by_type:
                by_type[etype] = []
            by_type[etype].append(ent)

        # Format context
        type_labels = {
            "preference": "Your preferences",
            "interest": "Your interests",
            "topic": "Topics you've discussed",
            "skill": "Your skills",
            "goal": "Your goals",
            "location": "Places",
            "organization": "Organizations",
            "person": "People you've mentioned",
            "fact": "Facts about you"
        }

        for etype, ents in by_type.items():
            label = type_labels.get(etype, etype.title())
            items = [f"- {e.name}" + (f": {e.description}" if e.description else "") for e in ents[:5]]
            if items:
                context_parts.append(f"\n**{label}:**")
                context_parts.extend(items)

        # Get some key relations
        relations = KnowledgeRelation.query.filter_by(user_id=user_id).order_by(
            KnowledgeRelation.strength.desc()
        ).limit(10).all()

        if relations:
            context_parts.append("\n**Key relationships:**")
            for rel in relations:
                source = rel.source_entity.name if rel.source_entity else "?"
                target = rel.target_entity.name if rel.target_entity else "?"
                if source == "User":
                    source = "You"
                context_parts.append(f"- {source} {rel.relation_type.replace('_', ' ')} {target}")

        return "\n".join(context_parts) if len(context_parts) > 1 else ""
    except Exception as e:
        print(f"[KNOWLEDGE] Failed to get user context: {str(e)}")
        return ""

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# Initialize OpenAI client for vision capabilities
openai_client = None
vision_enabled = False
vision_error = None

try:
    openai_api_key = os.environ.get('OPENAI_API_KEY')
    if openai_api_key:
        openai_client = OpenAI(api_key=openai_api_key)
        vision_enabled = True
        print("[SUCCESS] ✓ OpenAI Vision capability initialized")
    else:
        print("[WARNING] OPENAI_API_KEY not found - Vision capability disabled")
        vision_error = "OPENAI_API_KEY not configured"
except Exception as e:
    print(f"[ERROR] Failed to initialize OpenAI client: {str(e)}")
    vision_error = str(e)

def analyze_image_with_vision(filepath, filename, prompt=None):
    """
    Analyze an image using OpenAI's GPT-4 Vision API.
    Returns a detailed description of the image content.
    """
    start_time = time.time()

    if not openai_client:
        # Fallback to OCR if vision is not available
        try:
            image = Image.open(filepath)
            ocr_text = pytesseract.image_to_string(image)
            execution_time = int((time.time() - start_time) * 1000)
            if ocr_text.strip():
                result = f"[OCR Text Extracted - Vision not available]\n{ocr_text}"
                log_execution("OCR (Tesseract)", "vision", filename, result[:500], execution_time, "success",
                             extra_info={"fallback": True, "reason": "Vision API not configured"})
                return result
            log_execution("OCR (Tesseract)", "vision", filename, "No text detected", execution_time, "success",
                         extra_info={"fallback": True, "reason": "Vision API not configured"})
            return "[Image file - No text detected. Vision API not configured for image analysis]"
        except Exception as e:
            execution_time = int((time.time() - start_time) * 1000)
            log_execution("OCR (Tesseract)", "vision", filename, None, execution_time, "error", str(e))
            return f"[Image file - Could not process: {str(e)}]"

    try:
        # Read and encode the image
        with open(filepath, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')

        # Determine the image type
        file_ext = filename.rsplit('.', 1)[1].lower()
        mime_types = {
            'png': 'image/png',
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'gif': 'image/gif',
            'bmp': 'image/bmp',
            'webp': 'image/webp'
        }
        mime_type = mime_types.get(file_ext, 'image/jpeg')

        # Default prompt for image analysis
        if not prompt:
            prompt = """Analyze this image in detail. Provide:
1. A comprehensive description of what you see
2. Any text visible in the image (transcribe it)
3. Key objects, people, or elements
4. Colors, layout, and composition
5. Any relevant context or meaning

Be thorough but concise."""

        # Call OpenAI Vision API
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",  # Using gpt-4o-mini for cost efficiency, supports vision
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{base64_image}",
                                "detail": "high"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1500
        )

        vision_analysis = response.choices[0].message.content
        execution_time = int((time.time() - start_time) * 1000)

        # Also try OCR for any text that might be in the image
        try:
            image = Image.open(filepath)
            ocr_text = pytesseract.image_to_string(image)
            if ocr_text.strip():
                vision_analysis += f"\n\n[OCR Extracted Text]:\n{ocr_text.strip()}"
        except:
            pass  # OCR is optional, don't fail if it doesn't work

        # Log successful Vision API execution
        log_execution("GPT-4 Vision", "vision", filename, vision_analysis[:500], execution_time, "success",
                     extra_info={"model": "gpt-4o-mini", "mime_type": mime_type})

        return vision_analysis

    except Exception as e:
        execution_time = int((time.time() - start_time) * 1000)
        # Fallback to OCR on error
        try:
            image = Image.open(filepath)
            ocr_text = pytesseract.image_to_string(image)
            if ocr_text.strip():
                result = f"[Vision API error - OCR fallback]\n{ocr_text}"
                log_execution("OCR (Tesseract)", "vision", filename, result[:500], execution_time, "success",
                             extra_info={"fallback": True, "vision_error": str(e)})
                return result
            log_execution("GPT-4 Vision", "vision", filename, None, execution_time, "error", str(e))
            return f"[Image analysis failed: {str(e)}]"
        except Exception as ocr_error:
            log_execution("GPT-4 Vision", "vision", filename, None, execution_time, "error",
                         f"Vision: {str(e)}, OCR: {str(ocr_error)}")
            return f"[Image file - Could not analyze: {str(e)}]"

# Perplexica Search Tool - Open source AI-powered search engine
# GitHub: https://github.com/ItzCrazyKns/Perplexica

class SearchInput(BaseModel):
    """Input schema for the search tool."""
    query: str = Field(description="The search query to look up")

class PerplexicaSearchTool(BaseTool):
    """Search tool that uses Perplexica - an open source AI-powered search engine."""
    name: str = "Web Search"
    description: str = "Search the web for current information, news, and facts using AI-powered search. Use this when you need up-to-date information about any topic."
    args_schema: Type[BaseModel] = SearchInput
    backend_url: str = ""
    chat_model_provider: str = "openai"
    chat_model_name: str = "gpt-4o-mini"
    embedding_model_provider: str = "openai"
    embedding_model_name: str = "text-embedding-3-small"

    def __init__(self, backend_url: str = "", chat_model_provider: str = "openai",
                 chat_model_name: str = "gpt-4o-mini", embedding_model_provider: str = "openai",
                 embedding_model_name: str = "text-embedding-3-small", **kwargs):
        super().__init__(**kwargs)
        self.backend_url = backend_url
        self.chat_model_provider = chat_model_provider
        self.chat_model_name = chat_model_name
        self.embedding_model_provider = embedding_model_provider
        self.embedding_model_name = embedding_model_name

    def _get_focus_mode(self, query: str) -> str:
        """Determine the best focus mode based on the query."""
        query_lower = query.lower()

        # News-related queries
        if any(kw in query_lower for kw in ['news', 'latest', 'recent', 'today', 'breaking', 'update', 'happening']):
            return "webSearch"

        # Academic/research queries
        if any(kw in query_lower for kw in ['research', 'study', 'paper', 'scientific', 'academic', 'journal']):
            return "academicSearch"

        # Video queries
        if any(kw in query_lower for kw in ['video', 'youtube', 'watch', 'tutorial video']):
            return "youtubeSearch"

        # Reddit queries
        if any(kw in query_lower for kw in ['reddit', 'discussion', 'opinions', 'community']):
            return "redditSearch"

        # Default to web search
        return "webSearch"

    def _run(self, query: str) -> str:
        """Execute the search query using Perplexica."""
        start_time = time.time()
        focus_mode = None
        try:
            if not self.backend_url:
                log_execution("Perplexica", "search", query, None, 0, "error", "Backend URL not configured")
                return "Error: Perplexica backend URL not configured"

            # Determine focus mode based on query
            focus_mode = self._get_focus_mode(query)

            # Add current date context for news/recent queries to get up-to-date results
            search_query = query
            query_lower = query.lower()
            current_date = datetime.now()
            current_year = current_date.year
            current_month = current_date.strftime("%B %Y")

            # Check if query is asking for recent/current news
            news_keywords = ['news', 'latest', 'recent', 'today', 'current', 'now', 'update', 'happening', 'this week', 'this month']
            if any(kw in query_lower for kw in news_keywords):
                # Append current date context if not already present
                if str(current_year) not in query and current_date.strftime("%B") not in query:
                    search_query = f"{query} {current_month}"
                    print(f"[SEARCH] News query detected. Using date: {current_month}. Query: {search_query}")

            # Build request payload
            payload = {
                "chatModel": {
                    "provider": self.chat_model_provider,
                    "model": self.chat_model_name
                },
                "embeddingModel": {
                    "provider": self.embedding_model_provider,
                    "model": self.embedding_model_name
                },
                "focusMode": focus_mode,
                "query": search_query,
                "history": []
            }

            # Make request to Perplexica API
            response = requests.post(
                f"{self.backend_url}/api/search",
                json=payload,
                headers={"Content-Type": "application/json"},
                timeout=60
            )

            execution_time = int((time.time() - start_time) * 1000)

            if response.status_code != 200:
                error_msg = f"Perplexica returned status {response.status_code}"
                log_execution("Perplexica", "search", query, None, execution_time, "error", error_msg,
                             extra_info={"search_engine": "Perplexica", "focus_mode": focus_mode, "status_code": response.status_code})
                return f"Search error: {error_msg}"

            result = response.json()

            # Format the response
            formatted_output = []

            # Get the main answer/message
            if "message" in result:
                formatted_output.append(f"**Answer:**\n{result['message']}\n")

            # Get sources/citations if available
            sources_count = 0
            if "sources" in result and result["sources"]:
                sources_count = len(result["sources"])
                formatted_output.append("\n**Sources:**")
                for i, source in enumerate(result["sources"][:5], 1):
                    title = source.get("title", "Unknown")
                    url = source.get("url", "")
                    formatted_output.append(f"{i}. [{title}]({url})")

            output = "\n".join(formatted_output) if formatted_output else f"No results found for: {query}"

            # Log successful execution
            log_execution("Perplexica", "search", query, output[:500], execution_time, "success",
                         extra_info={"search_engine": "Perplexica", "focus_mode": focus_mode, "sources_count": sources_count})

            return output

        except requests.exceptions.Timeout:
            execution_time = int((time.time() - start_time) * 1000)
            log_execution("Perplexica", "search", query, None, execution_time, "timeout",
                         "Request timed out", extra_info={"search_engine": "Perplexica", "focus_mode": focus_mode})
            return "Search error: Request timed out. Perplexica server may be slow or unavailable."
        except requests.exceptions.ConnectionError:
            execution_time = int((time.time() - start_time) * 1000)
            log_execution("Perplexica", "search", query, None, execution_time, "error",
                         "Connection error", extra_info={"search_engine": "Perplexica", "focus_mode": focus_mode})
            return "Search error: Could not connect to Perplexica. Make sure it's running."
        except Exception as e:
            execution_time = int((time.time() - start_time) * 1000)
            log_execution("Perplexica", "search", query, None, execution_time, "error",
                         str(e), extra_info={"search_engine": "Perplexica", "focus_mode": focus_mode})
            return f"Search error: {str(e)}"

# DuckDuckGo Fallback Search Tool (free, no API key required)
class DuckDuckGoSearchTool(BaseTool):
    """Free search tool using DuckDuckGo - no API key required."""
    name: str = "Web Search"
    description: str = "Search the web for current information, news, and facts. Use this when you need up-to-date information about any topic."
    args_schema: Type[BaseModel] = SearchInput

    def _run(self, query: str) -> str:
        """Execute the search query using DuckDuckGo."""
        start_time = time.time()
        try:
            from duckduckgo_search import DDGS

            # Add current date for news queries
            search_query = query
            query_lower = query.lower()
            current_date = datetime.now()
            current_year = current_date.year
            current_month = current_date.strftime("%B %Y")

            # Check if query is asking for recent/current news
            news_keywords = ['news', 'latest', 'recent', 'today', 'current', 'now', 'update', 'happening']
            is_news = any(kw in query_lower for kw in news_keywords)

            if is_news and str(current_year) not in query:
                search_query = f"{query} {current_month}"
                print(f"[SEARCH] News query detected. Using date: {current_month}. Query: {search_query}")

            # Perform search
            with DDGS() as ddgs:
                if is_news:
                    # Use news search for news queries
                    results = list(ddgs.news(search_query, max_results=5))
                else:
                    # Use regular text search
                    results = list(ddgs.text(search_query, max_results=5))

            execution_time = int((time.time() - start_time) * 1000)

            if not results:
                log_execution("DuckDuckGo", "search", query, "No results", execution_time, "success",
                             extra_info={"search_engine": "DuckDuckGo", "is_news": is_news})
                return f"No results found for: {query}"

            # Format results
            output = []
            for i, r in enumerate(results, 1):
                title = r.get('title', 'No title')
                body = r.get('body', r.get('description', ''))
                url = r.get('href', r.get('url', r.get('link', '')))
                output.append(f"{i}. **{title}**\n   {body[:200]}...\n   🔗 {url}")

            result_text = "\n\n".join(output)

            log_execution("DuckDuckGo", "search", query, result_text[:500], execution_time, "success",
                         extra_info={"search_engine": "DuckDuckGo", "results_count": len(results), "is_news": is_news, "date_used": current_month, "actual_query": search_query})

            return result_text

        except Exception as e:
            execution_time = int((time.time() - start_time) * 1000)
            log_execution("DuckDuckGo", "search", query, None, execution_time, "error", str(e),
                         extra_info={"search_engine": "DuckDuckGo"})
            return f"Search error: {str(e)}"

# Initialize search tool - Perplexica primary, DuckDuckGo fallback
search_tool = None
search_tool_error = None
search_tool_type = None

# Try Perplexica first
perplexica_url = os.environ.get('PERPLEXICA_URL', '').strip().strip('"').strip("'")
if perplexica_url:
    print(f"[STARTUP] Initializing Perplexica search tool...")
    print(f"[INFO] Perplexica URL: {perplexica_url}")

    # Get model configuration from environment
    chat_provider = os.environ.get('PERPLEXICA_CHAT_PROVIDER', 'openai')
    chat_model = os.environ.get('PERPLEXICA_CHAT_MODEL', 'gpt-4o-mini')
    embed_provider = os.environ.get('PERPLEXICA_EMBED_PROVIDER', 'openai')
    embed_model = os.environ.get('PERPLEXICA_EMBED_MODEL', 'text-embedding-3-small')

    try:
        # Test connection to Perplexica
        test_response = requests.get(f"{perplexica_url}/api", timeout=10)
        if test_response.status_code == 200:
            search_tool = PerplexicaSearchTool(
                backend_url=perplexica_url,
                chat_model_provider=chat_provider,
                chat_model_name=chat_model,
                embedding_model_provider=embed_provider,
                embedding_model_name=embed_model
            )
            search_tool_type = "Perplexica"
            print(f"[SUCCESS] ✓ Perplexica search tool initialized")
            print(f"[INFO] Chat model: {chat_provider}/{chat_model}")
            print(f"[INFO] Embedding model: {embed_provider}/{embed_model}")
        else:
            print(f"[WARNING] Perplexica returned status {test_response.status_code}")
            search_tool_error = f"Perplexica returned status {test_response.status_code}"
    except requests.exceptions.ConnectionError:
        print(f"[WARNING] Could not connect to Perplexica at {perplexica_url}")
        search_tool_error = "Could not connect to Perplexica"
    except Exception as e:
        print(f"[WARNING] Perplexica initialization failed: {str(e)}")
        search_tool_error = str(e)
else:
    print(f"[INFO] PERPLEXICA_URL not set. Will use DuckDuckGo fallback.")

# Fall back to DuckDuckGo if Perplexica not available
if not search_tool:
    print(f"[STARTUP] Initializing DuckDuckGo search tool (free fallback)...")
    try:
        # Test if duckduckgo-search is available
        from duckduckgo_search import DDGS
        search_tool = DuckDuckGoSearchTool()
        search_tool_type = "DuckDuckGo"
        print(f"[SUCCESS] ✓ DuckDuckGo search tool initialized (no API key required)")
    except ImportError as e:
        print(f"[ERROR] DuckDuckGo search not available: {str(e)}")
        search_tool_error = "DuckDuckGo package not installed"
    except Exception as e:
        print(f"[ERROR] DuckDuckGo initialization failed: {str(e)}")
        search_tool_error = str(e)

if not search_tool:
    print(f"[WARNING] No search tool available. Install duckduckgo-search or set PERPLEXICA_URL.")

# Define the main conversational agent with search capabilities
# Only include search_tool if it was successfully initialized
agent_tools = [search_tool] if search_tool else []

conversational_agent = Agent(
    role="Intelligent AI Assistant",
    goal="""You are a versatile AI assistant similar to ChatGPT. Your goals are to:
    1. Engage in natural, helpful conversations
    2. Understand user intent and respond appropriately
    3. Use web search when current information is needed
    4. Reference uploaded documents when relevant
    5. Maintain conversation context and remember previous messages
    6. Ask for clarification when the user's request is ambiguous or unclear
    7. Be helpful, accurate, and conversational""",
    backstory="""You are an advanced AI assistant with multiple capabilities:
    - Natural conversation and general knowledge
    - Real-time web search for current information
    - Document analysis and retrieval
    - Context-aware responses that build on previous messages
    You intelligently decide when to search the web, when to use uploaded documents,
    and when to simply converse based on your knowledge. When uncertain about what
    the user wants, you ask clarifying questions.""",
    verbose=False,
    tools=agent_tools,
    allow_delegation=False
)

def calculate_file_hash(filepath):
    """Calculate MD5 hash of a file"""
    hash_md5 = hashlib.md5()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()

def format_response(text):
    """Format the response with better readability, convert URLs to hyperlinks, detect news widgets, and render tables"""

    # Check if response contains news articles (detect patterns)
    if "```news" in text or "```json" in text:
        # Extract JSON data for news widgets
        json_match = re.search(r'```(?:news|json)\s*(\[.*?\])\s*```', text, re.DOTALL)
        if json_match:
            try:
                news_data = json.loads(json_match.group(1))
                # Create news widget HTML
                news_html = '<div class="news-widgets">'
                for article in news_data[:6]:  # Max 6 articles
                    image = article.get('image', '')
                    title = article.get('title', 'No title')
                    description = article.get('description', '')
                    url = article.get('url', '#')
                    source = article.get('source', 'Unknown')

                    news_html += f'''
                    <div class="news-card">
                        {f'<img src="{image}" alt="{title}" class="news-image">' if image else '<div class="news-image-placeholder">📰</div>'}
                        <div class="news-content">
                            <h4 class="news-title">{title}</h4>
                            <p class="news-description">{description[:150]}...</p>
                            <div class="news-footer">
                                <span class="news-source">📌 {source}</span>
                                <a href="{url}" target="_blank" class="news-link">Read more →</a>
                            </div>
                        </div>
                    </div>
                    '''
                news_html += '</div>'
                # Remove the JSON from text and add widget
                text = re.sub(r'```(?:news|json)\s*\[.*?\]\s*```', news_html, text, flags=re.DOTALL)
            except:
                pass

    # Convert markdown tables to HTML tables
    def convert_markdown_table(match):
        table_text = match.group(0)
        lines = [line.strip() for line in table_text.split('\n') if line.strip()]

        if len(lines) < 2:
            return table_text

        # Parse header
        headers = [cell.strip() for cell in lines[0].split('|') if cell.strip()]

        # Skip separator line (line with dashes)
        data_lines = lines[2:] if len(lines) > 2 else []

        # Build HTML table
        html = '<div class="table-wrapper"><table class="formatted-table">'
        html += '<thead><tr>'
        for header in headers:
            html += f'<th>{header}</th>'
        html += '</tr></thead>'

        html += '<tbody>'
        for line in data_lines:
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            html += '<tr>'
            for cell in cells:
                html += f'<td>{cell}</td>'
            html += '</tr>'
        html += '</tbody></table></div>'

        return html

    # Match markdown tables (lines with | separators)
    table_pattern = r'(?:^|\n)(\|.+\|(?:\n\|.+\|)+)(?:\n|$)'
    text = re.sub(table_pattern, convert_markdown_table, text, flags=re.MULTILINE)

    # Convert URLs to HTML links
    url_pattern = r'(https?://[^\s<>"{}|\\^`\[\]]+)'
    text = re.sub(url_pattern, r'<a href="\1" target="_blank">\1</a>', text)

    # Convert **bold** to HTML bold
    text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)

    # Convert *italic* to HTML italic
    text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', text)

    # Convert line breaks to HTML breaks
    text = text.replace('\n\n', '<br><br>')
    text = text.replace('\n', '<br>')

    # Format numbered lists
    text = re.sub(r'^(\d+\.)', r'<br>\1', text, flags=re.MULTILINE)

    # Format bullet points
    text = re.sub(r'^- ', r'<br>• ', text, flags=re.MULTILINE)

    return text

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('home'))

    if request.method == 'POST':
        data = request.get_json()
        email = data.get('email')
        password = data.get('password')

        user = User.query.filter_by(email=email).first()
        if user and user.check_password(password):
            login_user(user)
            return jsonify({'success': True, 'message': 'Login successful!'})
        else:
            return jsonify({'success': False, 'message': 'Invalid email or password'}), 401

    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('home'))

    if request.method == 'POST':
        data = request.get_json()
        email = data.get('email')
        password = data.get('password')

        if User.query.filter_by(email=email).first():
            return jsonify({'success': False, 'message': 'Email already registered'}), 400

        user = User(email=email)
        user.set_password(password)
        db.session.add(user)
        db.session.commit()

        login_user(user)
        return jsonify({'success': True, 'message': 'Registration successful!'})

    return render_template('register.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

@app.route('/')
@login_required
def home():
    return render_template('index.html')

@app.route('/architecture')
@login_required
def architecture():
    return render_template('architecture.html')

@app.route('/logs')
@login_required
def logs():
    return render_template('logs.html')

# Knowledge Graph Routes
@app.route('/knowledge')
@login_required
def knowledge():
    return render_template('knowledge.html')

@app.route('/api/knowledge')
@login_required
def get_knowledge():
    """Get the user's knowledge graph data"""
    try:
        entities = KnowledgeEntity.query.filter_by(user_id=current_user.id).order_by(
            KnowledgeEntity.access_count.desc()
        ).all()

        relations = KnowledgeRelation.query.filter_by(user_id=current_user.id).all()

        entities_data = []
        for ent in entities:
            attrs = None
            if ent.attributes:
                try:
                    attrs = json.loads(ent.attributes)
                except:
                    attrs = ent.attributes

            entities_data.append({
                'id': ent.id,
                'type': ent.entity_type,
                'name': ent.name,
                'description': ent.description,
                'attributes': attrs,
                'confidence': ent.confidence,
                'source': ent.source,
                'access_count': ent.access_count,
                'created_at': ent.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                'updated_at': ent.updated_at.strftime('%Y-%m-%d %H:%M:%S')
            })

        relations_data = []
        for rel in relations:
            relations_data.append({
                'id': rel.id,
                'source_id': rel.source_id,
                'source_name': rel.source_entity.name if rel.source_entity else None,
                'target_id': rel.target_id,
                'target_name': rel.target_entity.name if rel.target_entity else None,
                'relation_type': rel.relation_type,
                'description': rel.description,
                'strength': rel.strength,
                'created_at': rel.created_at.strftime('%Y-%m-%d %H:%M:%S')
            })

        return jsonify({
            'success': True,
            'entities': entities_data,
            'relations': relations_data,
            'stats': {
                'total_entities': len(entities_data),
                'total_relations': len(relations_data),
                'entity_types': list(set(e['type'] for e in entities_data))
            }
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/knowledge/entity', methods=['POST'])
@login_required
def add_entity():
    """Manually add an entity to the knowledge graph"""
    try:
        data = request.json
        entity = store_entity(
            user_id=current_user.id,
            entity_type=data.get('type', 'fact'),
            name=data.get('name'),
            description=data.get('description'),
            attributes=data.get('attributes'),
            source='manual'
        )
        if entity:
            return jsonify({'success': True, 'entity_id': entity.id})
        return jsonify({'success': False, 'error': 'Failed to create entity'}), 400
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/knowledge/entity/<int:entity_id>', methods=['DELETE'])
@login_required
def delete_entity(entity_id):
    """Delete an entity from the knowledge graph"""
    try:
        entity = KnowledgeEntity.query.filter_by(id=entity_id, user_id=current_user.id).first()
        if not entity:
            return jsonify({'success': False, 'error': 'Entity not found'}), 404
        db.session.delete(entity)
        db.session.commit()
        return jsonify({'success': True})
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/knowledge/clear', methods=['POST'])
@login_required
def clear_knowledge():
    """Clear all knowledge graph data for the current user"""
    try:
        KnowledgeRelation.query.filter_by(user_id=current_user.id).delete()
        KnowledgeEntity.query.filter_by(user_id=current_user.id).delete()
        db.session.commit()
        return jsonify({'success': True, 'message': 'Knowledge graph cleared'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/logs')
@login_required
def get_logs():
    """API endpoint to fetch execution logs"""
    try:
        # Get query parameters
        limit = request.args.get('limit', 100, type=int)
        offset = request.args.get('offset', 0, type=int)
        tool_type = request.args.get('type', None)
        status = request.args.get('status', None)

        # Build query
        query = ExecutionLog.query.order_by(ExecutionLog.timestamp.desc())

        if tool_type:
            query = query.filter(ExecutionLog.tool_type == tool_type)
        if status:
            query = query.filter(ExecutionLog.status == status)

        # Get total count
        total = query.count()

        # Apply pagination
        logs = query.offset(offset).limit(limit).all()

        # Format logs for JSON response
        logs_data = []
        for log in logs:
            extra_info = None
            if log.extra_info:
                try:
                    extra_info = json.loads(log.extra_info)
                except:
                    extra_info = log.extra_info

            logs_data.append({
                'id': log.id,
                'timestamp': log.timestamp.strftime('%Y-%m-%d %H:%M:%S'),
                'tool_name': log.tool_name,
                'tool_type': log.tool_type,
                'input_data': log.input_data,
                'output_data': log.output_data,
                'execution_time_ms': log.execution_time_ms,
                'status': log.status,
                'error_message': log.error_message,
                'extra_info': extra_info
            })

        # Also include in-memory logs
        memory_logs = [{**log, 'source': 'memory'} for log in execution_logs_memory[-50:]]

        return jsonify({
            'success': True,
            'total': total,
            'logs': logs_data,
            'memory_logs': memory_logs
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/logs/stats')
@login_required
def get_log_stats():
    """Get statistics about execution logs"""
    try:
        total_executions = ExecutionLog.query.count()
        successful = ExecutionLog.query.filter_by(status='success').count()
        errors = ExecutionLog.query.filter_by(status='error').count()
        timeouts = ExecutionLog.query.filter_by(status='timeout').count()

        # Average execution time
        from sqlalchemy import func
        avg_time = db.session.query(func.avg(ExecutionLog.execution_time_ms)).scalar() or 0

        # Tool type breakdown
        tool_stats = db.session.query(
            ExecutionLog.tool_type,
            func.count(ExecutionLog.id)
        ).group_by(ExecutionLog.tool_type).all()

        # Tool name breakdown
        tool_name_stats = db.session.query(
            ExecutionLog.tool_name,
            func.count(ExecutionLog.id)
        ).group_by(ExecutionLog.tool_name).all()

        return jsonify({
            'success': True,
            'stats': {
                'total': total_executions,
                'successful': successful,
                'errors': errors,
                'timeouts': timeouts,
                'avg_execution_time_ms': round(avg_time, 2),
                'by_type': {t[0]: t[1] for t in tool_stats},
                'by_tool': {t[0]: t[1] for t in tool_name_stats}
            }
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/logs/clear', methods=['POST'])
@login_required
def clear_logs():
    """Clear all execution logs"""
    try:
        ExecutionLog.query.delete()
        db.session.commit()
        execution_logs_memory.clear()
        return jsonify({'success': True, 'message': 'Logs cleared'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

def extract_text_from_file(filepath, filename):
    """Extract text from different file types"""
    file_ext = filename.rsplit('.', 1)[1].lower()

    if file_ext == 'pdf':
        # Extract text from PDF
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    elif file_ext in ['docx', 'doc']:
        # Extract text from Word documents
        doc = Document(filepath)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    elif file_ext in ['pptx', 'ppt']:
        # Extract text from PowerPoint
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif file_ext in ['xlsx', 'xls']:
        # Extract text from Excel
        wb = load_workbook(filepath)
        text = ""
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text + "\n"
        return text

    elif file_ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp']:
        # Analyze images using Vision AI (with OCR fallback)
        return analyze_image_with_vision(filepath, filename)

    else:
        # Read text files (txt, md)
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

@app.route('/upload', methods=['POST'])
@login_required
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400

        file = request.files['file']
        session_id = request.form.get('session_id')

        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Calculate file hash
            file_hash = calculate_file_hash(filepath)

            # Check if file already exists for this user
            existing_doc = UploadedDocument.query.filter_by(file_hash=file_hash, user_id=current_user.id).first()
            if existing_doc:
                os.remove(filepath)  # Remove duplicate file
                return jsonify({
                    'success': True,
                    'message': f'📋 Document "{existing_doc.filename}" is already uploaded!',
                    'filename': existing_doc.filename,
                    'file_type': existing_doc.file_type,
                    'already_exists': True,
                    'char_count': len(existing_doc.content)
                })

            # Extract text/content from file (uses Vision AI for images)
            try:
                content = extract_text_from_file(filepath, filename)
            except Exception as extract_error:
                # Handle extraction errors gracefully
                file_ext = filename.rsplit('.', 1)[1].lower()
                if file_ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp']:
                    content = f"[Image file - Analysis not available: {str(extract_error)}]"
                else:
                    raise extract_error

            file_type = filename.rsplit('.', 1)[1].lower()

            # Store in database
            new_doc = UploadedDocument(
                user_id=current_user.id,
                filename=filename,
                file_hash=file_hash,
                content=content,
                file_type=file_type
            )
            db.session.add(new_doc)
            db.session.commit()

            # If session_id provided, add system message to chat
            if session_id:
                session = ChatSession.query.filter_by(id=session_id, user_id=current_user.id).first()
                if session:
                    system_msg = Message(
                        session_id=session_id,
                        role='system',
                        content=f'📎 Uploaded: {filename} ({file_type.upper()}, {len(content)} characters)'
                    )
                    db.session.add(system_msg)
                    db.session.commit()

            return jsonify({
                'success': True,
                'message': f'✅ File "{filename}" uploaded and indexed successfully!',
                'filename': filename,
                'file_type': file_type,
                'already_exists': False,
                'char_count': len(content)
            })
        else:
            return jsonify({'error': '❌ File type not allowed. Allowed types: txt, md, pdf, doc, docx, ppt, pptx, xls, xlsx, png, jpg, jpeg, gif, bmp'}), 400

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/documents', methods=['GET'])
@login_required
def get_documents():
    """Get all uploaded documents for current user"""
    try:
        docs = UploadedDocument.query.filter_by(user_id=current_user.id).order_by(UploadedDocument.uploaded_at.desc()).all()
        return jsonify({
            'documents': [{
                'id': doc.id,
                'filename': doc.filename,
                'file_type': doc.file_type,
                'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M')
            } for doc in docs]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/sessions', methods=['GET'])
@login_required
def get_sessions():
    """Get all chat sessions for current user"""
    try:
        sessions = ChatSession.query.filter_by(user_id=current_user.id).order_by(ChatSession.updated_at.desc()).all()
        return jsonify({
            'sessions': [{
                'id': session.id,
                'title': session.title,
                'created_at': session.created_at.strftime('%Y-%m-%d %H:%M'),
                'updated_at': session.updated_at.strftime('%Y-%m-%d %H:%M'),
                'message_count': len(session.messages)
            } for session in sessions]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/session/new', methods=['POST'])
@login_required
def new_session():
    """Create a new chat session for current user"""
    try:
        new_session = ChatSession(user_id=current_user.id, title="New Chat")
        db.session.add(new_session)
        db.session.commit()

        return jsonify({
            'success': True,
            'session_id': new_session.id,
            'title': new_session.title
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/session/<int:session_id>', methods=['GET'])
@login_required
def get_session(session_id):
    """Get a specific chat session with its messages"""
    try:
        session = ChatSession.query.filter_by(id=session_id, user_id=current_user.id).first_or_404()
        messages = Message.query.filter_by(session_id=session_id).order_by(Message.timestamp).all()

        return jsonify({
            'session': {
                'id': session.id,
                'title': session.title,
                'created_at': session.created_at.strftime('%Y-%m-%d %H:%M')
            },
            'messages': [{
                'role': msg.role,
                'content': msg.content,
                'timestamp': msg.timestamp.strftime('%Y-%m-%d %H:%M')
            } for msg in messages]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/session/<int:session_id>/delete', methods=['DELETE'])
@login_required
def delete_session(session_id):
    """Delete a chat session"""
    try:
        session = ChatSession.query.filter_by(id=session_id, user_id=current_user.id).first_or_404()
        db.session.delete(session)
        db.session.commit()

        return jsonify({'success': True, 'message': 'Session deleted'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint to verify configuration"""
    perplexica_url = os.environ.get('PERPLEXICA_URL', '')
    serp_key_raw = os.environ.get('SERPAPI_API_KEY', '')

    response = {
        'status': 'healthy',
        'search_enabled': search_tool is not None,
        'search_tool_type': search_tool_type,
        'perplexica_configured': bool(perplexica_url),
        'perplexica_url': perplexica_url if perplexica_url else None,
        'serpapi_configured': bool(serp_key_raw),
        'vision_enabled': vision_enabled,
        'openai_configured': bool(os.environ.get('OPENAI_API_KEY'))
    }
    if search_tool_error:
        response['search_error'] = search_tool_error
    if vision_error:
        response['vision_error'] = vision_error
    return jsonify(response)

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    try:
        user_message = request.json.get('message', '')
        use_rag = request.json.get('use_rag', False)
        session_id = request.json.get('session_id')
        agent_mode = request.json.get('mode', 'balanced')  # summarize, balanced, verbose, deep_thinking

        if not user_message:
            return jsonify({'error': 'No message provided'}), 400

        # Get or create session
        if not session_id:
            new_session = ChatSession(user_id=current_user.id, title="New Chat")
            db.session.add(new_session)
            db.session.commit()
            session_id = new_session.id

        session = ChatSession.query.filter_by(id=session_id, user_id=current_user.id).first()
        if not session:
            return jsonify({'error': 'Session not found'}), 404

        # Update session title if it's the first message
        if len(session.messages) == 0:
            # Generate title from first message
            title = user_message[:50] + "..." if len(user_message) > 50 else user_message
            session.title = title
            db.session.commit()

        # Add user message to database
        user_msg = Message(session_id=session_id, role='user', content=user_message)
        db.session.add(user_msg)
        db.session.commit()

        # Build conversation context from session history
        messages = Message.query.filter_by(session_id=session_id).order_by(Message.timestamp).all()
        context = ""
        if len(messages) > 1:
            # Include last 10 messages for context (5 exchanges)
            recent_history = messages[-11:-1] if len(messages) > 11 else messages[:-1]
            context = "Previous conversation:\n"
            for msg in recent_history:
                context += f"{msg.role.capitalize()}: {msg.content}\n"
            context += "\n"

        # Build document context if RAG is enabled (filter by current user)
        doc_context = ""
        if use_rag:
            docs = UploadedDocument.query.filter_by(user_id=current_user.id).all()
            if docs:
                doc_context = "\n\nAvailable documents:\n"
                for doc in docs:
                    # Limit document content to prevent token overflow
                    limited_content = doc.content[:2000] + "..." if len(doc.content) > 2000 else doc.content
                    doc_context += f"\n[{doc.filename}]:\n{limited_content}\n"

        # Get knowledge graph context about the user
        knowledge_context = get_user_knowledge_context(current_user.id, user_message)
        if knowledge_context:
            knowledge_context = f"\n\n{knowledge_context}\n"

        # Define mode-specific instructions
        mode_instructions = {
            'summarize': """- Provide concise, brief responses
- Focus on key points only
- Use bullet points when possible
- Keep responses short and to the point""",
            'balanced': """- Provide balanced, comprehensive responses
- Include relevant details without being overly verbose
- Use natural conversation flow
- Be helpful and concise""",
            'verbose': """- Provide detailed, thorough responses
- Explain concepts comprehensively
- Include examples and context
- Break down complex topics
- Use multiple paragraphs for clarity""",
            'deep_thinking': """- Analyze the question deeply from multiple angles
- Consider implications and nuances
- Provide reasoning and thought process
- Explore edge cases and alternatives
- Give comprehensive, well-reasoned responses
- Take time to think through the answer carefully"""
        }

        mode_instruction = mode_instructions.get(agent_mode, mode_instructions['balanced'])

        # Create intelligent task description
        task_description = f"""{context}Current user message: {user_message}
{doc_context}{knowledge_context}

You are an intelligent conversational AI assistant operating in **{agent_mode.replace('_', ' ').title()} Mode**.

Mode-specific guidelines:
{mode_instruction}

Analyze the user's message and respond appropriately:

1. **Simple Conversation**: If the user is greeting, asking a general question, or having casual conversation, respond naturally without using search tools.

2. **Document Query**: If the user asks about uploaded documents (and documents are available), use the document content to answer their question.

3. **Current Information/News**: If the user asks about recent events, news, current data, or anything requiring up-to-date information, use the search tool. When providing news:
   - If you find news articles with images and URLs, format them as a JSON array in markdown code block like this:
   ```news
   [
     {{"title": "Article Title", "description": "Brief description", "url": "https://...", "source": "Source Name", "image": "https://image-url.jpg"}},
     ...
   ]
   ```

4. **Ambiguous Requests**: If the user's request is unclear, vague, or could mean multiple things, politely ask for clarification instead of guessing. For example:
   - "I'm not sure I understand. Could you clarify what you mean by..."
   - "To help you better, could you provide more details about..."
   - "Are you asking about X or Y? Please let me know so I can assist you better."

5. **Context Awareness**: Remember and reference previous messages in the conversation when relevant.

Guidelines:
- Be conversational and natural with appropriate emojis
- Use search tool ONLY when current/factual information is needed
- Reference documents when the question relates to them
- Maintain context from previous messages
- If you use search, include source URLs
- When uncertain or confused, ASK for clarification
- Be helpful and concise"""

        # Create task
        task = Task(
            description=task_description,
            agent=conversational_agent,
            expected_output="A natural, context-aware response that appropriately addresses the user's message or asks for clarification if needed."
        )

        # Create and execute crew
        crew = Crew(
            agents=[conversational_agent],
            tasks=[task]
        )

        chat_start_time = time.time()
        try:
            result = crew.kickoff()
            chat_execution_time = int((time.time() - chat_start_time) * 1000)

            # Log successful chat execution
            log_execution(
                "CrewAI Agent", "agent",
                user_message[:500],
                str(result)[:1000],
                chat_execution_time,
                "success",
                extra_info={"mode": agent_mode, "use_rag": use_rag, "session_id": session_id}
            )
        except Exception as agent_error:
            chat_execution_time = int((time.time() - chat_start_time) * 1000)
            log_execution(
                "CrewAI Agent", "agent",
                user_message[:500],
                None,
                chat_execution_time,
                "error",
                str(agent_error),
                extra_info={"mode": agent_mode, "use_rag": use_rag}
            )
            raise agent_error

        # Add assistant response to database
        assistant_response = str(result)
        assistant_msg = Message(session_id=session_id, role='assistant', content=assistant_response)
        db.session.add(assistant_msg)

        # Update session timestamp
        session.updated_at = datetime.utcnow()
        db.session.commit()

        # Update knowledge graph with new information from conversation (async-like, non-blocking)
        try:
            update_knowledge_graph(current_user.id, user_message, assistant_response)
        except Exception as kg_error:
            print(f"[KNOWLEDGE] Non-critical error updating knowledge graph: {str(kg_error)}")

        # Format the response
        formatted_result = format_response(assistant_response)

        # Determine what actions were taken for progress simulation
        actions_taken = []
        if use_rag and docs:
            actions_taken.append('documents')
        if 'search' in assistant_response.lower() or 'http' in assistant_response.lower():
            actions_taken.append('web_search')
        actions_taken.append('thinking')

        return jsonify({
            'response': formatted_result,
            'session_id': session_id,
            'session_title': session.title,
            'mode': agent_mode,
            'actions': actions_taken
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # For Railway deployment, bind to 0.0.0.0 and use PORT from environment
    port = int(os.environ.get('PORT', 5000))
    app.run(debug=False, host='0.0.0.0', port=port)
